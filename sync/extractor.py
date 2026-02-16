import pdfplumber
import docx
import openpyxl
from pptx import Presentation


def extract_text(file_path: str, mime_type: str) -> str:
    """Text aus einer Datei extrahieren basierend auf MIME-Type."""
    extractors = {
        "application/pdf": _extract_pdf,
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": _extract_docx,
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": _extract_xlsx,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": _extract_pptx,
        # Google Docs werden als PDF exportiert
        "application/vnd.google-apps.document": _extract_pdf,
        "application/vnd.google-apps.spreadsheet": _extract_pdf,
        "application/vnd.google-apps.presentation": _extract_pdf,
    }

    extractor = extractors.get(mime_type)
    if not extractor:
        return ""

    try:
        return extractor(file_path)
    except Exception as e:
        print(f"Fehler beim Extrahieren von {file_path}: {e}")
        return ""


def _extract_pdf(path: str) -> str:
    text_parts = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)
    return "\n".join(text_parts)


def _extract_docx(path: str) -> str:
    doc = docx.Document(path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs)


def _extract_xlsx(path: str) -> str:
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    text_parts = []
    for sheet in wb.worksheets:
        text_parts.append(f"[Tabellenblatt: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                text_parts.append(" | ".join(cells))
    return "\n".join(text_parts)


def _extract_pptx(path: str) -> str:
    prs = Presentation(path)
    text_parts = []
    for i, slide in enumerate(prs.slides, 1):
        text_parts.append(f"[Folie {i}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        text_parts.append(text)
    return "\n".join(text_parts)
